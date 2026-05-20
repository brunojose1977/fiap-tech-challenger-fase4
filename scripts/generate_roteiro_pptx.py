#!/usr/bin/env python3
"""Gera apresentação .pptx a partir do roteiro de 15 minutos."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "docs" / "Roteiro de apresentação" / "Roteiro-Apresentacao-15min.pptx"
DIAGRAM_V2 = ROOT / "docs" / "Documentos de arquitetura" / "Diagrama de Arquitetura de Infraestrutura-v2.png"
INFOGRAFICO = ROOT / "docs" / "infograficos" / "Infografico-Projeto-YOLOv8-Pose.png"

# Paleta alinhada ao PDF do roteiro
NAVY = RGBColor(0x1A, 0x36, 0x5D)
BLUE = RGBColor(0x2C, 0x52, 0x82)
SLATE = RGBColor(0x2D, 0x37, 0x48)
MUTED = RGBColor(0x4A, 0x55, 0x68)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF7, 0xFA, 0xFC)


def _set_slide_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_header_bar(slide, title: str, subtitle: str | None = None) -> None:
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Inches(0),
        Inches(0),
        Inches(13.333),
        Inches(1.05),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.45), Inches(0.18), Inches(12.4), Inches(0.55))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = WHITE

    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.45), Inches(0.72), Inches(12.4), Inches(0.3))
        sp = sub.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(12)
        sp.font.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)


def _add_bullets(
    slide,
    items: list[str],
    *,
    left: float = 0.55,
    top: float = 1.35,
    width: float = 12.2,
    height: float = 5.8,
    font_size: int = 18,
    level_indent: bool = True,
) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 1 if level_indent and item.startswith("  ") else 0
        p.font.size = Pt(font_size - 2 if p.level else font_size)
        p.font.color.rgb = SLATE
        p.space_after = Pt(8)
        p.line_spacing = 1.15


def _add_quote(slide, text: str, *, top: float = 5.5) -> None:
    box = slide.shapes.add_textbox(Inches(0.7), Inches(top), Inches(11.9), Inches(1.2))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(14)
    p.font.italic = True
    p.font.color.rgb = MUTED


def _add_table_slide(
    prs: Presentation,
    title: str,
    headers: list[str],
    rows: list[list[str]],
    *,
    subtitle: str | None = None,
    col_widths: list[float] | None = None,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, LIGHT_BG)
    _add_header_bar(slide, title, subtitle)

    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(
        n_rows,
        n_cols,
        Inches(0.5),
        Inches(1.35),
        Inches(12.3),
        Inches(0.45 * n_rows + 0.3),
    )
    table = table_shape.table

    if col_widths:
        for idx, w in enumerate(col_widths):
            table.columns[idx].width = Inches(w)

    for c, header in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = header
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(11)
            p.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE

    for r, row in enumerate(rows, start=1):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = value
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.color.rgb = SLATE
            if r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF7, 0xFA, 0xFC)


def _add_image_slide(
    prs: Presentation,
    title: str,
    image_path: Path,
    *,
    subtitle: str | None = None,
    caption: str | None = None,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, LIGHT_BG)
    _add_header_bar(slide, title, subtitle)

    if image_path.is_file():
        slide.shapes.add_picture(
            str(image_path),
            Inches(0.6),
            Inches(1.25),
            width=Inches(12.1),
        )
    else:
        _add_bullets(slide, [f"[Imagem não encontrada: {image_path.name}]"], top=2.0)

    if caption:
        _add_quote(slide, caption, top=6.85)


def _title_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, NAVY)

    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.7), Inches(1.4))
    tp = title_box.text_frame.paragraphs[0]
    tp.text = "Pipeline YOLOv8 Pose + AWS"
    tp.font.size = Pt(40)
    tp.font.bold = True
    tp.font.color.rgb = WHITE
    tp.alignment = PP_ALIGN.CENTER

    sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.5), Inches(11.7), Inches(1.0))
    sp = sub_box.text_frame.paragraphs[0]
    sp.text = (
        "Solução 1: detecção de situações de violência contra a mulher em vídeo\n"
        "Roteiro de apresentação — 15 minutos"
    )
    sp.font.size = Pt(20)
    sp.font.color.rgb = RGBColor(0xCB, 0xD5, 0xE0)
    sp.alignment = PP_ALIGN.CENTER

    note = slide.shapes.add_textbox(Inches(0.8), Inches(5.8), Inches(11.7), Inches(0.5))
    np = note.text_frame.paragraphs[0]
    np.text = "yolo_v8_pose_estimation_code_project"
    np.font.size = Pt(12)
    np.font.color.rgb = MUTED
    np.alignment = PP_ALIGN.CENTER


def build_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    _title_slide(prs)

    _add_table_slide(
        prs,
        "Agenda — visão geral do tempo (~15 min)",
        ["Bloco", "Tópico", "Tempo"],
        [
            ["1", "Documento de arquitetura", "2 min"],
            ["2", "Estrutura do projeto (IaC, Docker/ECS, Python, testes, observabilidade)", "4 min"],
            ["3", "Segurança", "2 min"],
            ["4", "Esteiras CI/CD (GitHub + GitLab)", "2 min"],
            ["5", "Solução 1 — detecção em vídeo + demonstração", "5 min"],
        ],
        subtitle="Ritmo sugerido: ~120–130 palavras/minuto",
        col_widths=[1.0, 9.3, 2.0],
    )

    # --- Bloco 1 ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, LIGHT_BG)
    _add_header_bar(slide, "Bloco 1 — Documento de arquitetura", "0:00 – 2:00")
    _add_bullets(
        slide,
        [
            "Solução 1: análise de vídeo com estimativa de pose (YOLOv8) e infra AWS como código",
            "Documento de Arquitetura V2: fluxo dev → GitHub → Terraform → S3 → ECS Fargate → CloudWatch",
            "Versão 2 também documenta Transcribe + ChatGPT (áudio); foco desta apresentação: vídeo + pose",
            "",
            "Na tela: Documento de Arquitetura V2.pdf ou diagrama de infraestrutura V2",
        ],
        font_size=17,
    )
    _add_quote(
        slide,
        "«Referência única: quem interage, onde ficam os dados, como roda o processamento, "
        "rastreio e segurança (IAM por função, sem credenciais no código).»",
    )

    _add_image_slide(
        prs,
        "Arquitetura de infraestrutura (V2)",
        DIAGRAM_V2,
        subtitle="Bloco 1 — Documento de arquitetura",
        caption="docs/Documentos de arquitetura/Diagrama de Arquitetura de Infraestrutura-v2.png",
    )

    _add_table_slide(
        prs,
        "Stack tecnológico",
        ["Camada", "Tecnologias"],
        [
            ["Aplicação", "Python 3.11, Ultralytics YOLOv8 (pose), OpenCV, NumPy, boto3"],
            ["Container", "Docker (python:3.11-slim), FFmpeg e libs para OpenCV"],
            [
                "AWS",
                "S3, ECR, ECS Fargate, IAM (OIDC + task roles), CloudWatch Logs, VPC/subnets/IGW/SG",
            ],
            ["IaC", "Terraform ≥ 1.5"],
            [
                "CI/CD",
                "GitHub Actions (Ruff, Pytest, build Docker, validate Terraform, push ECR, RunTask)",
            ],
            ["Padrões", "12-factor; pacote src/yolo_violence_pipeline"],
        ],
        subtitle="Bloco 1",
        col_widths=[2.5, 9.8],
    )

    # --- Bloco 2 ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, LIGHT_BG)
    _add_header_bar(slide, "Bloco 2 — Estrutura do projeto", "2:00 – 6:00")
    _add_bullets(
        slide,
        [
            "Na tela: README, terraform/, Dockerfile, src/yolo_violence_pipeline/, tests/, .github/workflows/",
            "",
            "2.1 Terraform (~1 min 30 s) — terraform apply cria ambiente repetível (IaC)",
            "2.2 Docker → ECR → ECS Fargate (~1 min)",
            "2.3 Sequência Python (~1 min) — CLI, S3, YOLO pose, heurística, vídeo anotado",
            "2.4 Testes (~30 s) — Ruff + Pytest sem baixar pesos YOLO no CI",
            "2.5 CloudWatch Logs (~20 s) — stdout/stderr da task ECS",
        ],
        font_size=17,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, LIGHT_BG)
    _add_header_bar(slide, "2.1 Infraestrutura AWS via Terraform", "terraform/")
    _add_bullets(
        slide,
        [
            "VPC — default ou create_dedicated_vpc para VPC dedicada",
            "Subnets públicas (map_public_ip_on_launch); duas AZs na VPC dedicada",
            "Internet Gateway — rota 0.0.0.0/0 → IGW (Fargate baixa imagem do ECR)",
            "Tabelas de roteamento associadas às subnets públicas",
            "Security Group — egress para S3, ECR e APIs AWS",
            "Buckets S3 — entrada, predict, saída; Transcribe; bloqueio público, SSE-S3, versionamento",
        ],
        font_size=17,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, LIGHT_BG)
    _add_header_bar(slide, "2.2 Docker → ECR → ECS Fargate", "~1 min")
    _add_bullets(
        slide,
        [
            "Dockerfile: Python 3.11 + FFmpeg + dependências de vídeo + src/",
            "ECR: imagens versionadas",
            "ECS: cluster Fargate, task definition, execution role e task role",
            "CI: build + push; workflow RunTask com S3_INPUT_KEY",
            "",
            "Fluxo: git push → build → push ECR → RunTask → pipeline → S3 → CloudWatch",
        ],
        font_size=17,
    )

    _add_table_slide(
        prs,
        "2.3 Sequência de execução Python",
        ["Passo", "O que acontece"],
        [
            ["1", "CLI yolo-violence process carrega PipelineConfig do ambiente"],
            ["2", "boto3 baixa o vídeo de S3_INPUT_BUCKET + S3_INPUT_KEY"],
            ["3", "Ultralytics carrega o modelo (yolov8m-pose.pt ou yolov8n-pose.pt)"],
            ["4", "Predict task=pose, save=True → vídeo skeleton no bucket predict"],
            ["5", "Segunda passagem frame a frame com OpenCV + is_violence_detected"],
            ['6', 'Se disparar, overlay "Violence Detected!" no frame'],
            ["7", "MP4 final enviado ao bucket output"],
        ],
        subtitle="Heurística placeholder: keypoints visíveis acima de limiar",
        col_widths=[1.2, 11.1],
    )

    _add_table_slide(
        prs,
        "2.4 Cobertura de testes",
        ["Arquivo", "Foco"],
        [
            ["test_config.py", "Variáveis obrigatórias e defaults"],
            ["test_violence.py", "Heurística placeholder"],
            ["test_transcribe_*.py", "Módulos do fluxo Transcribe"],
        ],
        subtitle="ci-cd.yml: Ruff + Pytest em push/PR na branch master",
        col_widths=[4.5, 7.8],
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, LIGHT_BG)
    _add_header_bar(slide, "2.5 Observabilidade — CloudWatch Logs", "~20 s")
    _add_bullets(
        slide,
        [
            "Log group /ecs/.../yolo-violence — retenção 14 dias",
            "Mensagens: 'Violence no frame N', erros de S3",
            "Primeiro lugar para debugar RunTask em PENDING ou falha de download",
        ],
        font_size=20,
    )

    if INFOGRAFICO.is_file():
        _add_image_slide(
            prs,
            "Solução 1 — visão do pipeline",
            INFOGRAFICO,
            subtitle="Bloco 2 / contexto da demo",
        )

    # --- Bloco 3 ---
    _add_table_slide(
        prs,
        "Bloco 3 — Secrets no GitHub",
        ["Secret", "Uso"],
        [
            ["AWS_ROLE_ARN", "Assume role via OIDC (sem access key estática)"],
            ["ECR_REPOSITORY_NAME", "Push da imagem"],
            ["ECS_CLUSTER_NAME, ECS_TASK_DEFINITION_FAMILY, ECS_CONTAINER_NAME", "RunTask"],
            ["ECS_SUBNET_IDS, ECS_SECURITY_GROUP_ID", "Rede awsvpc"],
            ["OPENAI_API_KEY", "Apenas fluxo Transcribe (nunca no repositório)"],
        ],
        subtitle="6:00 – 8:00 | Na tela: Settings → Secrets (sem revelar valores)",
        col_widths=[5.5, 6.8],
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, LIGHT_BG)
    _add_header_bar(slide, "3.2 IAM e Security Groups", "~1 min 15 s")
    _add_bullets(
        slide,
        [
            "GitHub Actions (OIDC): token.actions.githubusercontent.com; trust repo:ORG/REPO:*",
            "ECS Execution Role: pull de imagem e logs",
            "ECS Task Role: políticas customizadas nos buckets e transcribe:*",
            "Security Group: apenas egress; sem portas de entrada desnecessárias",
            "",
            "Princípio: menor privilégio — CI não é admin; container não vê secrets do GitHub",
        ],
        font_size=17,
    )
    _add_quote(
        slide,
        "«Nenhuma chave AWS no código nem na imagem Docker. No ECS, credenciais vêm da task role.»",
        top=5.8,
    )

    # --- Bloco 4 ---
    _add_table_slide(
        prs,
        "Bloco 4 — GitHub Actions",
        ["Workflow", "Gatilho", "Função"],
        [
            ["ci-cd.yml", "Push/PR master + workflow_dispatch", "Ruff, Pytest, docker, terraform; push ECR"],
            ["run-fargate.yml", "Manual", "Push opcional + RunTask + validação rede"],
            ["aws-transcribe-...yml", "Manual", "Transcribe + ChatGPT (complementar)"],
        ],
        subtitle="8:00 – 10:00 | Mostrar: Actions verde, jobs lint-test/docker/terraform/publish-ecr",
        col_widths=[3.8, 4.2, 4.3],
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, LIGHT_BG)
    _add_header_bar(slide, "GitLab — equivalente conceitual", "~1 min")
    _add_bullets(
        slide,
        [
            "Pipelines versionados no GitHub neste repositório",
            "Equivalente GitLab: .gitlab-ci.yml com lint, test, docker build, terraform validate, OIDC AWS",
            "Mostrar captura institucional se exigido pela banca",
            "",
            "Nota: não há .gitlab-ci.yml aqui — mesmo desenho CI, outro orquestrador",
        ],
        font_size=17,
    )

    # --- Bloco 5 ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, LIGHT_BG)
    _add_header_bar(slide, "Bloco 5 — Solução 1: violência em vídeo", "10:00 – 15:00")
    _add_bullets(
        slide,
        [
            "YOLOv8-pose no S3: YOLOv8m-pose (precisão) ou YOLOv8n-pose (demo/CI leve)",
            "OpenCV remonta vídeo anotado; boto3 download/upload; Python 3.11 no Fargate ou local",
            "",
            "Stack: entrada S3 → predict (pose) → keypoints → is_violence_detected → saída S3 + predict",
        ],
        font_size=17,
    )

    _add_table_slide(
        prs,
        "Roteiro da demonstração",
        ["Tempo", "Ação", "O que mostrar"],
        [
            ["10:00–11:00", "Vídeo original", "originais/original-video01.mp4"],
            ["11:00–11:45", "Objeto no S3", "bucket *-input-* → meu_video_04.mp4"],
            ["11:45–12:30", "Processamento", "Slide fluxo 7 passos ou log CloudWatch"],
            ["12:30–13:30", "Vídeo processado", "processado/video01_violence_detected_*.mp4"],
            ["13:30–14:15", "Saída no S3", "bucket *-output-* → video_violence_detected_*.mp4"],
            ["14:15–15:00", "Fechamento", "Limitações éticas/técnicas + próximos passos"],
        ],
        subtitle="Bloco 5 — demo ao vivo",
        col_widths=[2.2, 3.0, 7.1],
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, LIGHT_BG)
    _add_header_bar(slide, "Falas para a demo", "Bloco 5")
    _add_bullets(
        slide,
        [
            "Original: material bruto — gravação que chegaria ao bucket de entrada",
            "S3 entrada: Fargate lê por chave S3 — desacopla upload do processamento",
            "Processado: keypoints + regra sinaliza frames suspeitos — revisão humana",
            "S3 saída: resultado versionado, pronto para dashboard ou perícia",
        ],
        font_size=18,
        top=1.4,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, NAVY)
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.7), Inches(2.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Fechamento"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    body = tf.add_paragraph()
    body.text = (
        "Arquitetura documentada • Infra como código • Container no ECR • "
        "Fargate serverless • CI/CD com testes e OIDC • Pose + vídeo anotado no S3"
    )
    body.font.size = Pt(18)
    body.font.color.rgb = RGBColor(0xCB, 0xD5, 0xE0)
    body.alignment = PP_ALIGN.CENTER
    body.space_before = Pt(24)

    next_p = tf.add_paragraph()
    next_p.text = (
        "Próximo passo: modelo supervisionado no lugar do placeholder; "
        "YOLOv8m-pose vs latência no Fargate CPU"
    )
    next_p.font.size = Pt(16)
    next_p.font.color.rgb = MUTED
    next_p.alignment = PP_ALIGN.CENTER
    next_p.space_before = Pt(20)

    thanks = tf.add_paragraph()
    thanks.text = "Obrigado. Perguntas?"
    thanks.font.size = Pt(28)
    thanks.font.bold = True
    thanks.font.color.rgb = WHITE
    thanks.alignment = PP_ALIGN.CENTER
    thanks.space_before = Pt(36)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, LIGHT_BG)
    _add_header_bar(slide, "Checklist pré-apresentação", "Speaker notes")
    _add_bullets(
        slide,
        [
            "☐ PDF/PNG de arquitetura V2 aberto",
            "☐ Console AWS: buckets input / predict / output",
            "☐ Vídeos local: originais/ e processado/",
            "☐ GitHub Actions: pipeline verde + run-fargate",
            "☐ GitLab: URL/captura institucional (se exigido)",
            "☐ CloudWatch: log stream de execução bem-sucedida",
            "☐ Não exibir valores de secrets na tela",
            "",
            "Modelo: default yolov8n-pose.pt no ECS; use MODEL_NAME=yolov8m-pose.pt para maior precisão",
        ],
        font_size=17,
    )

    return prs


def main() -> None:
    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs = build_presentation()
    prs.save(str(OUTPUT))
    print(f"PPTX gerado: {OUTPUT} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
